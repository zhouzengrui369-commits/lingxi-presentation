"""PPTX → JSON 提取器（python-pptx 驱动）。

灵犀演示 · Phase 1 · T-1.3
被 cli_template.ts / cli.ts 通过 subprocess 调用：
    python backend/scripts/extract_pptx.py <input.pptx> > <output.json>

输出 JSON 结构：
{
  "file": "...",
  "slide_width_emu": 12192000,
  "slide_height_emu": 6858000,
  "slides": [
    {
      "index": 0,
      "layout_type_guess": "title",
      "shapes": [
        {"type": "textbox", "left": ..., "top": ..., "width": ..., "height": ...,
         "text": "2026 战略规划", "runs": [{"text": ..., "font_name": "...",
         "font_size_pt": 44, "bold": true, "color_rgb": "0B1F3A"}]},
        {"type": "rect", "fill_rgb": "0B1F3A", ...},
      ]
    }
  ],
  "extracted_at": "2026-07-09T..."
}
"""
from __future__ import annotations

import json
import sys
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# Force UTF-8 stdout on Windows (cp1252 default breaks Chinese chars in PPTX text)
if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8")
if hasattr(sys.stderr, "reconfigure"):
    sys.stderr.reconfigure(encoding="utf-8")


def _shape_bounding_emu(shape) -> tuple[int, int, int, int]:
    return (
        int(shape.left or 0),
        int(shape.top or 0),
        int(shape.width or 0),
        int(shape.height or 0),
    )


def _color_to_hex(color_format) -> str | None:
    """从 pptx ColorFormat 取 RGB hex（无颜色返回 None）。"""
    try:
        if color_format.type is None:
            return None
        rgb = color_format.rgb
        if rgb is None:
            return None
        return str(rgb).upper()
    except (AttributeError, ValueError):
        return None


def _extract_run(run) -> dict[str, Any]:
    out: dict[str, Any] = {
        "text": run.text,
        "font_name": None,
        "font_size_pt": None,
        "bold": None,
        "italic": None,
        "color_rgb": None,
    }
    try:
        if run.font.name:
            out["font_name"] = run.font.name
    except (AttributeError, ValueError):
        pass
    try:
        if run.font.size:
            out["font_size_pt"] = int(run.font.size.pt)
    except (AttributeError, ValueError, TypeError):
        pass
    try:
        if run.font.bold is not None:
            out["bold"] = bool(run.font.bold)
    except (AttributeError, ValueError):
        pass
    try:
        if run.font.italic is not None:
            out["italic"] = bool(run.font.italic)
    except (AttributeError, ValueError):
        pass
    color_hex = _color_to_hex(run.font.color)
    if color_hex:
        out["color_rgb"] = color_hex
    return out


def _extract_textbox(shape) -> dict[str, Any]:
    left, top, width, height = _shape_bounding_emu(shape)
    info: dict[str, Any] = {
        "type": "textbox",
        "left": left,
        "top": top,
        "width": width,
        "height": height,
        "text": "",
        "runs": [],
    }
    tf = shape.text_frame
    parts: list[str] = []
    for para in tf.paragraphs:
        for run in para.runs:
            info["runs"].append(_extract_run(run))
            parts.append(run.text)
    info["text"] = " ".join(parts).strip()
    return info


def _extract_rect_or_shape(shape) -> dict[str, Any]:
    left, top, width, height = _shape_bounding_emu(shape)
    info: dict[str, Any] = {
        "type": "rect",
        "left": left,
        "top": top,
        "width": width,
        "height": height,
        "fill_rgb": None,
        "line_rgb": None,
        "shape_name": shape.name,
    }
    fill_hex = _color_to_hex(shape.fill.fore_color)
    if fill_hex:
        info["fill_rgb"] = fill_hex
    try:
        line_hex = _color_to_hex(shape.line.color)
        if line_hex:
            info["line_rgb"] = line_hex
    except (AttributeError, ValueError, TypeError):
        pass
    return info


def _guess_layout_type(slide_index: int, total: int, shapes: list[dict[str, Any]]) -> str:
    """启发式版式推断（基于形状数 + 文本分布 + 位置）。"""
    text_shapes = [s for s in shapes if s["type"] == "textbox"]
    rect_shapes = [s for s in shapes if s["type"] == "rect"]

    n_text = len(text_shapes)
    n_rect = len(rect_shapes)

    if n_text == 0:
        return "blank"

    # 第一页 / 最后一页通常 title
    if slide_index == 0:
        return "title"
    if slide_index == total - 1:
        # 看是否还有文本
        if any(s.get("text") for s in text_shapes):
            return "summary"
        return "blank"

    # 看文本框数：2 个 = two_column / quote
    if n_text == 2:
        # 看是否有一个超大字（font > 30）
        sizes = []
        for s in text_shapes:
            for r in s.get("runs", []):
                if r.get("font_size_pt"):
                    sizes.append(r["font_size_pt"])
        if sizes and max(sizes) >= 28:
            return "quote"
        return "two_column"

    # 看是否有大标题（font > 24）+ 副标题（font 14-20）
    big_titles = 0
    for s in text_shapes:
        for r in s.get("runs", []):
            if r.get("font_size_pt") and r["font_size_pt"] >= 24 and r.get("bold"):
                big_titles += 1
    if big_titles == 1 and n_text == 2:
        return "section"

    if big_titles >= 1:
        return "content"

    return "content"


def extract_pptx(pptx_path: Path) -> dict[str, Any]:
    """主入口：解析 PPTX → 结构化 JSON。"""
    pres = Presentation(str(pptx_path))
    total = len(pres.slides)
    slides_json: list[dict[str, Any]] = []

    for idx, slide in enumerate(pres.slides):
        shapes_json: list[dict[str, Any]] = []
        for shape in slide.shapes:
            shape_type = shape.shape_type
            try:
                if shape.has_text_frame:
                    shapes_json.append(_extract_textbox(shape))
                elif shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE or shape_type == MSO_SHAPE_TYPE.FREEFORM:
                    shapes_json.append(_extract_rect_or_shape(shape))
                elif shape_type == MSO_SHAPE_TYPE.PICTURE:
                    left, top, width, height = _shape_bounding_emu(shape)
                    shapes_json.append({
                        "type": "picture",
                        "left": left,
                        "top": top,
                        "width": width,
                        "height": height,
                    })
                else:
                    # 兜底
                    left, top, width, height = _shape_bounding_emu(shape)
                    shapes_json.append({
                        "type": "other",
                        "left": left,
                        "top": top,
                        "width": width,
                        "height": height,
                        "shape_type": str(shape_type),
                    })
            except Exception as exc:  # noqa: BLE001
                shapes_json.append({
                    "type": "error",
                    "error": str(exc),
                    "shape_name": getattr(shape, "name", "?"),
                })

        slide_json = {
            "index": idx,
            "layout_type_guess": _guess_layout_type(idx, total, shapes_json),
            "shapes": shapes_json,
        }
        slides_json.append(slide_json)

    return {
        "file": pptx_path.name,
        "file_path": str(pptx_path),
        "slide_width_emu": int(pres.slide_width),
        "slide_height_emu": int(pres.slide_height),
        "slide_count": total,
        "slides": slides_json,
        "extracted_at": datetime.now(timezone.utc).isoformat(),
        "extractor_version": "1.0.0",
    }


def main(argv: list[str]) -> int:
    if len(argv) != 2:
        print(f"Usage: {argv[0]} <input.pptx>", file=sys.stderr)
        return 2
    pptx_path = Path(argv[1])
    if not pptx_path.exists():
        print(f"Error: file not found: {pptx_path}", file=sys.stderr)
        return 1
    data = extract_pptx(pptx_path)
    json.dump(data, sys.stdout, ensure_ascii=False, indent=2)
    sys.stdout.write("\n")
    return 0


if __name__ == "__main__":
    sys.exit(main(sys.argv))