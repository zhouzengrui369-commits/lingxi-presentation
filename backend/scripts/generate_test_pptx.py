"""生成 3 套测试 PPTX 模板（业务深色 / 学术浅色 / 创意渐变）。

每套 ≥ 10 页，每页有不同的版式 + 配色 + 字体，便于 T-1.3 style_analyzer 验证。

灵犀演示 · Phase 1 · T-1.3
"""
from __future__ import annotations

import sys
from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


# ---- 工具：给段落加 run ----

def _add_run(paragraph, text: str, *, font_name: str | None = None,
             font_size: int | None = None, bold: bool = False,
             color_rgb: RGBColor | None = None) -> None:
    run = paragraph.add_run()
    run.text = text
    if font_name:
        run.font.name = font_name
    if font_size:
        run.font.size = Pt(font_size)
    if bold:
        run.font.bold = True
    if color_rgb:
        run.font.color.rgb = color_rgb


def _add_text_box(slide, text: str, *, left: int, top: int, width: int, height: int,
                  font_name: str | None = None, font_size: int = 18,
                  bold: bool = False, color_rgb: RGBColor | None = None,
                  align: PP_ALIGN = PP_ALIGN.LEFT,
                  anchor: MSO_ANCHOR = MSO_ANCHOR.TOP) -> None:
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    _add_run(p, text, font_name=font_name, font_size=font_size, bold=bold, color_rgb=color_rgb)


def _add_filled_rect(slide, *, left: int, top: int, width: int, height: int,
                     fill_rgb: RGBColor, line_rgb: RGBColor | None = None) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if line_rgb is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_rgb
    shape.text_frame.text = ""  # 不显示文字


# ---- 主题 1：商务深色（business-dark） ----

def build_business_dark(out: Path) -> None:
    """12 页商务深色模板：深蓝底 + 白色文字 + 金色强调。"""
    pres = Presentation()
    pres.slide_width = Inches(13.333)
    pres.slide_height = Inches(7.5)

    bg = RGBColor(0x0B, 0x1F, 0x3A)        # 深海军蓝
    accent = RGBColor(0xE0, 0xB3, 0x4F)     # 金色
    text_white = RGBColor(0xFF, 0xFF, 0xFF)
    text_gray = RGBColor(0xB0, 0xB8, 0xC9)
    body_font = "Microsoft YaHei"
    title_font = "Microsoft YaHei"

    # 12 页：title / section / content / two_column / quote / summary / blank 等版式
    layouts_content = [
        ("title",       "2026 战略规划",       "商务增长部 · 年度汇报"),
        ("section",     "一、市场回顾",       "Section 01"),
        ("content",     "核心业务增长",       "全年营收同比增长 18%，利润率提升 3.2 个百分点"),
        ("two_column",  "国内 vs 海外",       "国内稳健 / 海外加速"),
        ("content",     "客户结构升级",       "Top 10 客户贡献占比从 35% 提升至 48%"),
        ("quote",       "增长来自深耕",       "—— 首席战略官"),
        ("content",     "二、挑战与机遇",     "供应链 / 政策 / 技术三重变革"),
        ("two_column",  "风险 vs 机会",       "识别 4 大风险 / 抓住 3 大机会"),
        ("content",     "三、2026 目标",      "营收 +25%，新增 3 条业务线"),
        ("content",     "执行路线图",         "Q1 基础 / Q2 突破 / Q3 验证 / Q4 复制"),
        ("summary",     "总结",               "聚焦核心客户 + 加速海外 + 技术投入"),
        ("blank",       "谢谢聆听",           "Q&A"),
    ]

    blank_layout = pres.slide_layouts[6]  # Blank layout
    for i, (layout_type, heading, body) in enumerate(layouts_content):
        slide = pres.slides.add_slide(blank_layout)

        # 全背景深色填充
        _add_filled_rect(slide, left=0, top=0, width=13.333, height=7.5, fill_rgb=bg)

        # 顶部装饰条
        if layout_type in ("title", "section", "summary"):
            _add_filled_rect(slide, left=0, top=0, width=13.333, height=0.12,
                             fill_rgb=accent)
            _add_filled_rect(slide, left=0, top=7.38, width=13.333, height=0.12,
                             fill_rgb=accent)

        if layout_type == "title":
            # 居中大标题
            _add_text_box(slide, heading, left=1, top=2.5, width=11.3, height=1.2,
                          font_name=title_font, font_size=44, bold=True,
                          color_rgb=text_white, align=PP_ALIGN.CENTER)
            _add_text_box(slide, body, left=1, top=4.0, width=11.3, height=0.8,
                          font_name=body_font, font_size=20, color_rgb=text_gray,
                          align=PP_ALIGN.CENTER)
        elif layout_type == "section":
            _add_text_box(slide, "01", left=0.5, top=2.0, width=2, height=1.5,
                          font_name=title_font, font_size=72, bold=True,
                          color_rgb=accent)
            _add_text_box(slide, heading, left=2.5, top=2.5, width=10, height=1.0,
                          font_name=title_font, font_size=36, bold=True,
                          color_rgb=text_white)
            _add_text_box(slide, body, left=2.5, top=3.7, width=10, height=0.6,
                          font_name=body_font, font_size=16, color_rgb=text_gray)
        elif layout_type == "two_column":
            # 左列
            _add_text_box(slide, heading, left=0.5, top=0.5, width=12.3, height=0.8,
                          font_name=title_font, font_size=28, bold=True,
                          color_rgb=text_white)
            # 中间分割线
            _add_filled_rect(slide, left=6.6, top=2.0, width=0.02, height=5.0,
                             fill_rgb=accent)
            _add_text_box(slide, "国内业务", left=0.5, top=2.0, width=6, height=0.6,
                          font_name=body_font, font_size=22, bold=True,
                          color_rgb=accent)
            _add_text_box(slide, "稳健增长，季度环比 +5%", left=0.5, top=2.8,
                          width=6, height=1.5, font_name=body_font,
                          font_size=18, color_rgb=text_white)
            _add_text_box(slide, "海外业务", left=7.0, top=2.0, width=6, height=0.6,
                          font_name=body_font, font_size=22, bold=True,
                          color_rgb=accent)
            _add_text_box(slide, "高速增长，季度环比 +18%", left=7.0, top=2.8,
                          width=6, height=1.5, font_name=body_font,
                          font_size=18, color_rgb=text_white)
        elif layout_type == "quote":
            _add_text_box(slide, '"', left=0.5, top=1.5, width=1.5, height=1.5,
                          font_name=title_font, font_size=80, bold=True,
                          color_rgb=accent)
            _add_text_box(slide, heading, left=1.8, top=2.5, width=10, height=1.5,
                          font_name=title_font, font_size=32, bold=True,
                          color_rgb=text_white)
            _add_text_box(slide, body, left=1.8, top=4.5, width=10, height=0.6,
                          font_name=body_font, font_size=18, color_rgb=text_gray)
        else:
            # content / summary / blank — 标题 + 正文
            if layout_type != "blank":
                _add_text_box(slide, heading, left=0.5, top=0.5, width=12.3, height=0.8,
                              font_name=title_font, font_size=28, bold=True,
                              color_rgb=text_white)
                # 短装饰线
                _add_filled_rect(slide, left=0.5, top=1.35, width=0.8, height=0.04,
                                 fill_rgb=accent)
                _add_text_box(slide, body, left=0.5, top=2.0, width=12.3, height=4.5,
                              font_name=body_font, font_size=20, color_rgb=text_white)
            else:
                # blank
                _add_text_box(slide, heading, left=1, top=3.0, width=11.3, height=1.5,
                              font_name=title_font, font_size=60, bold=True,
                              color_rgb=text_white, align=PP_ALIGN.CENTER)
                _add_text_box(slide, body, left=1, top=5.0, width=11.3, height=0.6,
                              font_name=body_font, font_size=20, color_rgb=accent,
                              align=PP_ALIGN.CENTER)

    pres.save(out)
    print(f"  [OK] {out.name}: {len(pres.slides)} pages")


# ---- 主题 2：学术浅色（academic-light） ----

def build_academic_light(out: Path) -> None:
    """12 页学术浅色模板：白底 + 深灰文字 + 蓝色标题，衬线字体。"""
    pres = Presentation()
    pres.slide_width = Inches(13.333)
    pres.slide_height = Inches(7.5)

    bg = RGBColor(0xFA, 0xFA, 0xF7)        # 米白
    accent = RGBColor(0x1F, 0x4E, 0x79)     # 深学术蓝
    text_dark = RGBColor(0x22, 0x22, 0x22)
    text_muted = RGBColor(0x55, 0x55, 0x55)
    rule = RGBColor(0xCC, 0xCC, 0xCC)
    body_font = "SimSun"
    title_font = "SimHei"

    layouts_content = [
        ("title",       "深度学习综述",            "从感知到推理的演进"),
        ("section",     "第一章 引言",            "Section 1"),
        ("content",     "研究背景",                "深度学习已成为 AI 主流范式"),
        ("content",     "研究方法",                "本文采用系统性文献综述方法"),
        ("two_column",  "贡献 vs 局限",            "贡献 / 局限"),
        ("content",     "实验设置",                "ImageNet / GLUE / 9 项 benchmark"),
        ("quote",       "数据为王",                "—— 经典 AI 格言"),
        ("content",     "结果分析",                "在 7/9 任务上达到 SOTA"),
        ("two_column",  "对比 vs 展望",            "对比 / 展望"),
        ("content",     "讨论",                    "可解释性 / 鲁棒性 / 公平性"),
        ("summary",     "结论",                    "深度学习仍处于快速发展期"),
        ("blank",       "参考文献",                "References"),
    ]

    blank_layout = pres.slide_layouts[6]
    for i, (layout_type, heading, body) in enumerate(layouts_content):
        slide = pres.slides.add_slide(blank_layout)

        # 全背景米白填充
        _add_filled_rect(slide, left=0, top=0, width=13.333, height=7.5, fill_rgb=bg)

        if layout_type in ("title", "section", "summary"):
            _add_filled_rect(slide, left=0, top=0, width=0.15, height=7.5,
                             fill_rgb=accent)

        if layout_type == "title":
            _add_text_box(slide, heading, left=0.8, top=2.5, width=11.5, height=1.2,
                          font_name=title_font, font_size=44, bold=True,
                          color_rgb=accent, align=PP_ALIGN.CENTER)
            _add_filled_rect(slide, left=5.5, top=4.0, width=2.3, height=0.04,
                             fill_rgb=accent)
            _add_text_box(slide, body, left=0.8, top=4.3, width=11.5, height=0.8,
                          font_name=body_font, font_size=20, color_rgb=text_muted,
                          align=PP_ALIGN.CENTER)
        elif layout_type == "section":
            _add_text_box(slide, heading, left=0.8, top=3.0, width=11.5, height=1.2,
                          font_name=title_font, font_size=40, bold=True,
                          color_rgb=text_dark)
            _add_text_box(slide, body, left=0.8, top=4.5, width=11.5, height=0.6,
                          font_name=body_font, font_size=16, color_rgb=text_muted)
        elif layout_type == "two_column":
            _add_text_box(slide, heading, left=0.8, top=0.5, width=11.5, height=0.8,
                          font_name=title_font, font_size=28, bold=True,
                          color_rgb=text_dark)
            _add_filled_rect(slide, left=0.8, top=1.35, width=11.5, height=0.02,
                             fill_rgb=rule)
            _add_text_box(slide, "贡献", left=0.8, top=2.0, width=5.5, height=0.6,
                          font_name=title_font, font_size=22, bold=True,
                          color_rgb=accent)
            _add_text_box(slide, "提出统一框架", left=0.8, top=2.8,
                          width=5.5, height=1.5, font_name=body_font,
                          font_size=18, color_rgb=text_dark)
            _add_text_box(slide, "局限", left=7.0, top=2.0, width=5.5, height=0.6,
                          font_name=title_font, font_size=22, bold=True,
                          color_rgb=accent)
            _add_text_box(slide, "依赖大规模标注数据", left=7.0, top=2.8,
                          width=5.5, height=1.5, font_name=body_font,
                          font_size=18, color_rgb=text_dark)
        elif layout_type == "quote":
            _add_text_box(slide, '"', left=0.8, top=1.5, width=1.5, height=1.5,
                          font_name=title_font, font_size=80, bold=True,
                          color_rgb=accent)
            _add_text_box(slide, heading, left=2.2, top=2.5, width=10, height=1.5,
                          font_name=title_font, font_size=32, bold=True,
                          color_rgb=text_dark)
            _add_text_box(slide, body, left=2.2, top=4.5, width=10, height=0.6,
                          font_name=body_font, font_size=18, color_rgb=text_muted)
        else:
            if layout_type != "blank":
                _add_text_box(slide, heading, left=0.8, top=0.5, width=11.5, height=0.8,
                              font_name=title_font, font_size=28, bold=True,
                              color_rgb=text_dark)
                _add_filled_rect(slide, left=0.8, top=1.35, width=11.5, height=0.02,
                                 fill_rgb=rule)
                _add_text_box(slide, body, left=0.8, top=2.0, width=11.5, height=4.5,
                              font_name=body_font, font_size=20, color_rgb=text_dark)
            else:
                _add_text_box(slide, heading, left=1, top=3.0, width=11.3, height=1.5,
                              font_name=title_font, font_size=60, bold=True,
                              color_rgb=text_dark, align=PP_ALIGN.CENTER)
                _add_text_box(slide, body, left=1, top=5.0, width=11.3, height=0.6,
                              font_name=body_font, font_size=20, color_rgb=text_muted,
                              align=PP_ALIGN.CENTER)

    pres.save(out)
    print(f"  [OK] {out.name}: {len(pres.slides)} pages")


# ---- 主题 3：创意渐变（creative-gradient） ----

def build_creative_gradient(out: Path) -> None:
    """12 页创意渐变模板：紫粉橙渐变 + 白底文字，几何形状装饰。"""
    pres = Presentation()
    pres.slide_width = Inches(13.333)
    pres.slide_height = Inches(7.5)

    bg_top = RGBColor(0x6A, 0x11, 0xCB)     # 紫
    bg_mid = RGBColor(0xE9, 0x1E, 0x63)     # 粉
    bg_bot = RGBColor(0xFF, 0xC1, 0x07)     # 橙
    text_white = RGBColor(0xFF, 0xFF, 0xFF)
    accent = RGBColor(0xFF, 0xEB, 0x3B)     # 亮黄
    body_font = "STHeiti"
    title_font = "STHeiti"

    layouts_content = [
        ("title",       "创意设计 2026",          "让灵感变成现实"),
        ("section",     "PART 01 视觉",           "Section"),
        ("content",     "色彩理论",                "色彩是情感的第一语言"),
        ("two_column",  "形式 vs 功能",            "形式 / 功能"),
        ("content",     "材质与质感",              "材质让设计有温度"),
        ("quote",       "设计即态度",              "—— 创意总监"),
        ("content",     "动效原则",                "动效让界面有生命"),
        ("content",     "字体选择",                "字体是品牌的灵魂"),
        ("two_column",  "抽象 vs 具象",            "抽象 / 具象"),
        ("content",     "案例研究",                "客户满意度 +42%"),
        ("summary",     "下一步",                  "持续实验 + 用户共创"),
        ("blank",       "感谢观看",                "Thank You"),
    ]

    blank_layout = pres.slide_layouts[6]
    for i, (layout_type, heading, body) in enumerate(layouts_content):
        slide = pres.slides.add_slide(blank_layout)

        # 用三个色块模拟渐变背景
        _add_filled_rect(slide, left=0, top=0, width=13.333, height=2.5, fill_rgb=bg_top)
        _add_filled_rect(slide, left=0, top=2.5, width=13.333, height=2.5, fill_rgb=bg_mid)
        _add_filled_rect(slide, left=0, top=5.0, width=13.333, height=2.5, fill_rgb=bg_bot)

        # 装饰几何形
        if layout_type in ("title", "section", "summary"):
            _add_filled_rect(slide, left=10.5, top=0.5, width=2.3, height=2.3,
                             fill_rgb=accent)

        if layout_type == "title":
            _add_text_box(slide, heading, left=1, top=2.0, width=11.3, height=1.5,
                          font_name=title_font, font_size=48, bold=True,
                          color_rgb=text_white, align=PP_ALIGN.CENTER)
            _add_text_box(slide, body, left=1, top=4.0, width=11.3, height=0.8,
                          font_name=body_font, font_size=22, color_rgb=text_white,
                          align=PP_ALIGN.CENTER)
        elif layout_type == "section":
            _add_text_box(slide, heading, left=0.8, top=3.0, width=11.5, height=1.2,
                          font_name=title_font, font_size=40, bold=True,
                          color_rgb=text_white, align=PP_ALIGN.CENTER)
            _add_text_box(slide, body, left=0.8, top=4.5, width=11.5, height=0.6,
                          font_name=body_font, font_size=18, color_rgb=accent,
                          align=PP_ALIGN.CENTER)
        elif layout_type == "two_column":
            _add_text_box(slide, heading, left=0.5, top=0.3, width=12.3, height=0.8,
                          font_name=title_font, font_size=28, bold=True,
                          color_rgb=text_white)
            # 装饰形
            _add_filled_rect(slide, left=6.5, top=2.5, width=0.5, height=3.0,
                             fill_rgb=accent)
            _add_text_box(slide, "形式", left=0.8, top=2.5, width=5, height=0.6,
                          font_name=title_font, font_size=24, bold=True,
                          color_rgb=text_white)
            _add_text_box(slide, "形态即品牌", left=0.8, top=3.3,
                          width=5, height=1.5, font_name=body_font,
                          font_size=18, color_rgb=text_white)
            _add_text_box(slide, "功能", left=7.5, top=2.5, width=5, height=0.6,
                          font_name=title_font, font_size=24, bold=True,
                          color_rgb=accent)
            _add_text_box(slide, "好用才是硬道理", left=7.5, top=3.3,
                          width=5, height=1.5, font_name=body_font,
                          font_size=18, color_rgb=text_white)
        elif layout_type == "quote":
            _add_text_box(slide, heading, left=1.5, top=2.5, width=10, height=1.5,
                          font_name=title_font, font_size=36, bold=True,
                          color_rgb=text_white, align=PP_ALIGN.CENTER)
            _add_text_box(slide, body, left=1.5, top=4.5, width=10, height=0.6,
                          font_name=body_font, font_size=20, color_rgb=accent,
                          align=PP_ALIGN.CENTER)
        else:
            if layout_type != "blank":
                _add_text_box(slide, heading, left=0.5, top=0.3, width=12.3, height=0.8,
                              font_name=title_font, font_size=28, bold=True,
                              color_rgb=text_white)
                _add_text_box(slide, body, left=0.8, top=2.5, width=11.5, height=3.5,
                              font_name=body_font, font_size=20, color_rgb=text_white)
            else:
                _add_text_box(slide, heading, left=1, top=3.0, width=11.3, height=1.5,
                              font_name=title_font, font_size=60, bold=True,
                              color_rgb=text_white, align=PP_ALIGN.CENTER)
                _add_text_box(slide, body, left=1, top=5.0, width=11.3, height=0.6,
                              font_name=body_font, font_size=22, color_rgb=accent,
                              align=PP_ALIGN.CENTER)

    pres.save(out)
    print(f"  [OK] {out.name}: {len(pres.slides)} pages")


def main() -> int:
    out_dir = Path(__file__).resolve().parents[2] / "apps" / "desktop" / "testdata" / "templates"
    out_dir.mkdir(parents=True, exist_ok=True)

    print(f"[INFO] Generating 3 PPTX test fixtures in {out_dir}")
    build_business_dark(out_dir / "business-dark.pptx")
    build_academic_light(out_dir / "academic-light.pptx")
    build_creative_gradient(out_dir / "creative-gradient.pptx")
    print(f"[OK] 3 templates generated")
    return 0


if __name__ == "__main__":
    sys.exit(main())