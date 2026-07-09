#!/usr/bin/env python3
"""生成 T-1.1 7 格式样本 + 50MB+ 大文件压力测试样本。

用法：
    python scripts/generate_testdata.py
    python scripts/generate_testdata.py --output-dir /custom/path

输出：
    apps/desktop/testdata/
    ├── docx_sample.docx
    ├── pdf_sample.pdf
    ├── xlsx_sample.xlsx
    ├── pptx_sample.pptx
    ├── md_sample.md
    ├── jpg_sample.jpg
    ├── png_sample.png
    └── large_50mb.pdf        # ≥ 50MB 压测样本

灵犀演示 · Phase 1 · T-1.1
"""

from __future__ import annotations

import argparse
import io
import os
import sys
import zlib
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont

from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches, Pt
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont


SAMPLE_TEXT = """Q1 2026 季度汇报关键业绩

一、营收概况
Q1 营收达成率 108%，同比增长 23%。新增企业客户 14 家，流失 3 家。
AI 产品线贡献占比从 28% 提升至 41%，成为第一大营收支柱。

二、海外市场
海外市场初步开拓，东南亚 3 个 POC 客户进入商务谈判阶段。

三、产品里程碑
- AI 顾问引导闭环上线
- HTML 预览可控生成
- 多格式输出（PPTX / PDF / DOCX / HTML）

四、团队扩张
团队从 12 人扩张到 18 人，新引入 2 名高级 AI 工程师。

五、风险与挑战
- LLM Wiki 整理在大文件（>50M）下处理慢或失败
- React Native 桌面端（macOS/Win）生态不成熟，组件库选型踩坑
- PPTX 模板 → HTML 转换保真度不足（版式/字体丢失）
"""


def write_md(path: Path) -> None:
    """Markdown 样本 — 季度汇报纯文本。"""
    content = f"""# 灵犀演示 · 季度汇报 (Markdown 样本)

{SAMPLE_TEXT}

## 关键数据表

| 指标 | Q4 2025 | Q1 2026 | 同比 |
|---|---|---|---|
| 营收达成率 | 95% | 108% | +13pp |
| 同比增长 | 18% | 23% | +5pp |
| 新增客户 | 8 | 14 | +75% |

> 灵犀演示 · Phase 1 · T-1.1 文件管理 · LLM Wiki 自动整理
"""
    path.write_text(content, encoding="utf-8")


def write_docx(path: Path) -> None:
    """Word 样本 — 用 python-docx 生成。"""
    doc = Document()
    doc.add_heading("灵犀演示 · 季度汇报 (DOCX 样本)", level=1)
    for para in SAMPLE_TEXT.split("\n\n"):
        doc.add_paragraph(para.strip())
    doc.add_heading("关键数据", level=2)
    table = doc.add_table(rows=4, cols=4)
    table.style = "Light Grid Accent 1"
    headers = ["指标", "Q4 2025", "Q1 2026", "同比"]
    for i, h in enumerate(headers):
        table.rows[0].cells[i].text = h
    data = [
        ("营收达成率", "95%", "108%", "+13pp"),
        ("同比增长", "18%", "23%", "+5pp"),
        ("新增客户", "8", "14", "+75%"),
    ]
    for r, row in enumerate(data, 1):
        for c, v in enumerate(row):
            table.rows[r].cells[c].text = v
    doc.save(path)


def write_xlsx(path: Path) -> None:
    """Excel 样本 — 用 openpyxl 生成。"""
    wb = Workbook()
    ws = wb.active
    ws.title = "Q1 2026 业绩"
    ws["A1"] = "灵犀演示 · 季度汇报 (XLSX 样本)"
    ws["A3"] = "指标"
    ws["B3"] = "Q4 2025"
    ws["C3"] = "Q1 2026"
    ws["D3"] = "同比"
    rows = [
        ("营收达成率", "95%", "108%", "+13pp"),
        ("同比增长", "18%", "23%", "+5pp"),
        ("新增客户", 8, 14, "+75%"),
        ("AI 产品线占比", "28%", "41%", "+13pp"),
        ("团队规模", 12, 18, "+50%"),
    ]
    for r, row in enumerate(rows, 4):
        for c, v in enumerate(row, 1):
            ws.cell(row=r, column=c, value=v)
    # 第 2 张表：风险登记
    ws2 = wb.create_sheet("风险登记")
    ws2["A1"] = "风险"
    ws2["B1"] = "可能性"
    ws2["C1"] = "影响"
    ws2["D1"] = "缓解"
    risks = [
        ("LLM Wiki 大文件处理慢", "中", "高", "分块 + 进度条"),
        ("RN 桌面端组件缺失", "中", "中", "react-native-macos 兜底"),
        ("PPTX 转 HTML 保真度差", "中", "高", "python-pptx + 手工解析双路"),
    ]
    for r, row in enumerate(risks, 2):
        for c, v in enumerate(row, 1):
            ws2.cell(row=r, column=c, value=v)
    wb.save(path)


def write_pptx(path: Path) -> None:
    """PPTX 样本 — 5 页幻灯片。"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    titles = [
        ("灵犀演示 · 季度汇报", "AI 驱动的办公内容生成桌面 App"),
        ("营收概况", "Q1 营收达成率 108%，同比增长 23%"),
        ("产品里程碑", "AI 顾问引导 · HTML 预览 · 多格式输出"),
        ("海外市场", "东南亚 3 个 POC 客户进入商务谈判"),
        ("风险与展望", "LLM Wiki 大文件 + RN 桌面端 + PPTX 保真"),
    ]
    for title, body in titles:
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # title + content
        slide.shapes.title.text = title
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = body
        else:
            slide.shapes.add_textbox(Inches(1), Inches(2), Inches(10), Inches(4)).text = body

    # 第 6 页：纯文本 page 用作 fallback
    blank = prs.slides.add_slide(prs.slide_layouts[6])
    blank.shapes.add_textbox(Inches(1), Inches(1), Inches(11), Inches(5)).text = "谢谢\nQ&A"
    prs.save(path)


def write_pdf(path: Path, *, large: bool = False) -> None:
    """PDF 样本 — 用 reportlab 生成。"""
    c = canvas.Canvas(str(path), pagesize=letter)
    width, height = letter

    if large:
        # 大文件：每页放一段「高位熵」文本（伪随机但确定），
        # 关掉 reportlab 自身的 page compression，让 PDF 真正大。
        # 1200 页 × 每页 ~50KB ≈ 60MB
        c.pageCompression = 0
        import hashlib
        pages = 4000
        for i in range(pages):
            seed = hashlib.sha256(f"lingxi-page-{i}".encode()).digest()
            # 每页写 50KB unique data
            payload = bytearray()
            for j in range(50 * 1024):
                # 用 SHA256 链生成 deterministic random bytes
                if j % 32 == 0:
                    seed = hashlib.sha256(seed).digest()
                payload.append(seed[j % 32])
            # 把 payload 转成可读字符（4 chars per 3 bytes base64-ish）
            # 直接当 latin-1 字符串写
            text = "".join(chr(b % 95 + 32) for b in payload[: 45 * 1024])
            c.setFont("Helvetica", 6)
            c.drawString(20, height - 20, f"=== Page {i+1} ===")
            # 分行绘制
            y = height - 32
            line_len = 130
            for k in range(0, len(text), line_len):
                c.drawString(20, y, text[k : k + line_len])
                y -= 8
                if y < 20:
                    break
            c.showPage()
    else:
        c.setFont("Helvetica-Bold", 18)
        c.drawString(40, height - 60, "Lingxi Demo - Quarterly Review (PDF Sample)")
        c.setFont("Helvetica", 11)
        y = height - 100
        for para in SAMPLE_TEXT.split("\n\n"):
            for line in para.strip().split("\n"):
                c.drawString(40, y, line[:90])
                y -= 14
                if y < 40:
                    c.showPage()
                    y = height - 40
            y -= 8
        c.showPage()

    c.save()


def _png_chunk(chunk_type: bytes, data: bytes) -> bytes:
    """构造一个 PNG chunk（4B length + 4B type + data + 4B CRC）。"""
    import struct
    import zlib as _zl
    length = struct.pack(">I", len(data))
    crc = struct.pack(">I", _zl.crc32(chunk_type + data) & 0xFFFFFFFF)
    return length + chunk_type + data + crc


def write_jpg(path: Path) -> None:
    """JPG 样本 — PIL 生成 1024x768 「灵犀演示 · 季度汇报」图。"""
    img = Image.new("RGB", (1024, 768), color=(245, 247, 250))
    draw = ImageDraw.Draw(img)
    # 主标题
    try:
        # 尝试常见中文字体路径
        font_paths = [
            "/System/Library/Fonts/PingFang.ttc",
            "/System/Library/Fonts/STHeiti Light.ttc",
            "/System/Library/Fonts/Hiragino Sans GB.ttc",
            "/Library/Fonts/Arial Unicode.ttf",
        ]
        font_title = None
        font_body = None
        for fp in font_paths:
            if os.path.isfile(fp):
                font_title = ImageFont.truetype(fp, 48)
                font_body = ImageFont.truetype(fp, 24)
                break
        if not font_title:
            font_title = ImageFont.load_default()
            font_body = ImageFont.load_default()
    except Exception:
        font_title = ImageFont.load_default()
        font_body = ImageFont.load_default()

    draw.rectangle([(0, 0), (1024, 100)], fill=(30, 64, 175))
    draw.text((40, 25), "灵犀演示 · 季度汇报", fill=(255, 255, 255), font=font_title)
    draw.text((40, 140), "Q1 2026 营收达成率 108%", fill=(15, 23, 42), font=font_body)
    draw.text((40, 180), "同比增长 23%（AI 产品线占比 41%）", fill=(71, 85, 105), font=font_body)
    draw.text((40, 220), "新增企业客户 14 家", fill=(71, 85, 105), font=font_body)
    draw.text((40, 260), "海外市场 POC 3 家进入商务谈判", fill=(71, 85, 105), font=font_body)
    # 底部装饰
    draw.rectangle([(0, 700), (1024, 768)], fill=(226, 232, 240))
    draw.text((40, 720), "灵犀演示 · Phase 1 · T-1.1 文件管理 · LLM Wiki", fill=(100, 116, 139), font=font_body)
    img.save(path, "JPEG", quality=85)


def write_png(path: Path) -> None:
    """PNG 样本 — 透明背景 + 业务图标。"""
    img = Image.new("RGBA", (800, 600), color=(255, 255, 255, 255))
    draw = ImageDraw.Draw(img)
    try:
        font_paths = [
            "/System/Library/Fonts/PingFang.ttc",
            "/System/Library/Fonts/STHeiti Light.ttc",
            "/System/Library/Fonts/Hiragino Sans GB.ttc",
            "/Library/Fonts/Arial Unicode.ttf",
        ]
        font_big = None
        for fp in font_paths:
            if os.path.isfile(fp):
                font_big = ImageFont.truetype(fp, 40)
                break
        if not font_big:
            font_big = ImageFont.load_default()
    except Exception:
        font_big = ImageFont.load_default()

    draw.rectangle([(0, 0), (800, 80)], fill=(15, 23, 42))
    draw.text((30, 20), "Lingxi KB Architecture (PNG Sample)", fill=(255, 255, 255), font=font_big)
    # 简单流程图
    boxes = [
        ("Importer", 60, 200, (59, 130, 246)),
        ("Wiki", 320, 200, (16, 185, 129)),
        ("Storage", 580, 200, (245, 158, 11)),
        ("KB Linker", 200, 420, (139, 92, 246)),
        ("Advisor (T-1.2)", 460, 420, (236, 72, 153)),
    ]
    for label, x, y, color in boxes:
        draw.rectangle([(x, y), (x + 180, y + 80)], fill=color, outline=(15, 23, 42), width=2)
        draw.text((x + 20, y + 25), label, fill=(255, 255, 255), font=font_big)
    img.save(path, "PNG")


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("--output-dir", default=None, help="覆盖默认输出目录")
    args = parser.parse_args()

    out = Path(args.output_dir) if args.output_dir else (
        Path(__file__).resolve().parents[1] / "apps" / "desktop" / "testdata"
    )
    out.mkdir(parents=True, exist_ok=True)

    print(f"[testdata] output dir: {out}")

    targets = [
        ("md_sample.md", lambda p: write_md(p)),
        ("docx_sample.docx", lambda p: write_docx(p)),
        ("xlsx_sample.xlsx", lambda p: write_xlsx(p)),
        ("pptx_sample.pptx", lambda p: write_pptx(p)),
        ("pdf_sample.pdf", lambda p: write_pdf(p)),
        ("jpg_sample.jpg", lambda p: write_jpg(p)),
        ("png_sample.png", lambda p: write_png(p)),
    ]

    for name, fn in targets:
        path = out / name
        fn(path)
        size = path.stat().st_size
        print(f"  ✓ {name:30s}  {size:>12,d} bytes")

    # 大文件压测样本（≥ 50MB PDF）
    large = out / "large_50mb.pdf"
    print(f"[testdata] generating large PDF (this may take ~30s)...")
    write_pdf(large, large=True)
    size = large.stat().st_size
    print(f"  ✓ {'large_50mb.pdf':30s}  {size:>12,d} bytes ({size/1024/1024:.1f} MB)")

    return 0


if __name__ == "__main__":
    sys.exit(main())