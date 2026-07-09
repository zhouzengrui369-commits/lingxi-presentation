#!/usr/bin/env python3
"""生成 T-1.1 测试数据 — 7 格式样本 + 1 个 >50M 大 PDF

输出到 apps/desktop/testdata/：
  - sample.docx       Word 文档（季度汇报主题）
  - sample.pdf        PDF（同一主题，3 页）
  - sample.xlsx       Excel 表格（季度数据）
  - sample.pptx       PowerPoint（5 页幻灯片）
  - sample.md         Markdown 笔记
  - sample.jpg        JPG 图片（800x600, 加文字水印）
  - sample.png        PNG 图片（800x600, 加文字水印）
  - large_50mb.pdf    60MB+ PDF（100M 内压测样本）
  - README.md         testdata 索引

灵犀演示 · Phase 1 · T-1.1
"""
import os
import sys
from pathlib import Path

# Force UTF-8 stdout (CJK print)
import io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")

from docx import Document
from docx.shared import Pt
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches, Pt as PptxPt
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.cidfonts import UnicodeCIDFont
from PIL import Image, ImageDraw, ImageFont

# Register CJK font for reportlab
pdfmetrics.registerFont(UnicodeCIDFont("STSong-Light"))


HERE = Path(__file__).parent.resolve()
# testdata 目录与 scripts 同级（在 apps/desktop/testdata/）
OUT = HERE.parent / "testdata"
OUT.mkdir(parents=True, exist_ok=True)

# ---- 主题内容（季度汇报场景） ----

REPORT_TITLE = "Q1 2026 季度汇报 · 灵犀演示"

REPORT_PARAGRAPHS = [
    "Q1 2026 季度业绩总结：本季度营收达成率 108%，同比增长 23%。",
    "新增企业客户 14 家，流失 3 家，净增 11 家，续约率 92%。",
    "AI 产品线贡献占比从 28% 提升至 41%，成为第一大营收支柱。",
    "海外市场初步开拓，东南亚 3 个 POC 客户进入商务谈判阶段。",
    "团队规模 28 人，新增 5 人（研发 3 / 销售 2），流失 1 人。",
    "产品迭代：发布 6 个主要版本，平均迭代周期 14 天，线上事故 0 起。",
    "Q2 计划：聚焦 AI 商业化闭环，推出付费版 + 行业模板，目标 ARR 800 万。",
]

TAGS = ["季度汇报", "营收", "AI产品线", "海外市场", "团队", "Q2计划"]


def make_docx():
    """Word 样本"""
    doc = Document()
    title = doc.add_heading(REPORT_TITLE, level=1)
    for p in REPORT_PARAGRAPHS:
        doc.add_paragraph(p)
    doc.add_paragraph("").add_run("关键词：").bold = True
    doc.add_paragraph(" · ".join(TAGS))
    p = doc.add_paragraph()
    p.add_run("\n本文件由 T-1.1 testdata generator 生成。").italic = True
    out = OUT / "sample.docx"
    doc.save(out)
    print(f"  [OK] {out.name} ({out.stat().st_size} bytes)")


def make_pdf(pages: int = 3, target_size_bytes: int = 0, name: str = "sample.pdf"):
    """PDF 样本（可指定页数 + 目标大小）"""
    out = OUT / name

    if target_size_bytes > 0:
        # 大文件模式：先写最小 PDF，然后追加 padding（PDF 注释流）
        c = canvas.Canvas(str(out), pagesize=letter)
        c.setFont("STSong-Light", 22)
        c.drawString(72, 720, REPORT_TITLE)
        c.setFont("STSong-Light", 14)
        c.drawString(72, 680, "T-1.1 large file stress test (PRD ≥99% success rate @ 100M)")
        c.showPage()
        c.save()

        # 追加 raw bytes（PDF 阅读器忽略尾部垃圾；importer 只关心头部 metadata）
        size = out.stat().st_size
        need = target_size_bytes - size
        if need > 0:
            chunk = b"\n%PDF-PADDING-1MB\n" + (b" " * (1024 * 1024 - 32))
            with open(out, "ab") as f:
                written = 0
                while written < need:
                    f.write(chunk)
                    written += len(chunk)
        size = out.stat().st_size
        print(f"  [OK] {name} ({size} bytes = {size // (1024*1024)} MB)")
    else:
        # 正常 PDF
        c = canvas.Canvas(str(out), pagesize=letter)
        c.setFont("STSong-Light", 18)
        width, height = letter

        for page_idx in range(pages):
            c.setFont("STSong-Light", 22)
            c.drawString(72, height - 100, REPORT_TITLE)
            c.setFont("STSong-Light", 14)
            y = height - 160
            for i, para in enumerate(REPORT_PARAGRAPHS):
                c.drawString(72, y, f"  • {para}")
                y -= 24
                if y < 100:
                    c.showPage()
                    y = height - 80
                    c.setFont("STSong-Light", 14)
            c.setFont("STSong-Light", 11)
            c.drawString(72, 60, f"第 {page_idx + 1} / {pages} 页  ·  关键词: {', '.join(TAGS)}")
            c.showPage()

        c.save()
        size = out.stat().st_size
        print(f"  [OK] {name} ({size} bytes)")


def make_xlsx():
    """Excel 样本"""
    wb = Workbook()
    ws = wb.active
    ws.title = "Q1业绩"
    ws.append(["指标", "Q1", "Q4 2025", "同比增长", "备注"])
    rows = [
        ["营收(万)", 1280, 1040, "23%", "达成率 108%"],
        ["新增客户", 14, 9, "56%", "企业客户为主"],
        ["AI 产品线营收(万)", 525, 290, "81%", "占比 41%"],
        ["海外营收(万)", 95, 0, "—", "东南亚 3 POC"],
        ["续约率", "92%", "85%", "+7pp", "—"],
        ["团队人数", 28, 24, "+4", "研发 +3 销售 +2"],
    ]
    for r in rows:
        ws.append(r)
    ws2 = wb.create_sheet("Q2计划")
    ws2.append(["项目", "负责人", "目标", "完成时间"])
    for r in [
        ["付费版上线", "张三", "ARR 200 万", "2026-06-30"],
        ["行业模板", "李四", "5 套", "2026-05-31"],
        ["海外 POC 转化", "王五", "2 单", "2026-06-15"],
    ]:
        ws2.append(r)
    out = OUT / "sample.xlsx"
    wb.save(out)
    print(f"  [OK] {out.name} ({out.stat().st_size} bytes)")


def make_pptx():
    """PPTX 样本（5 页）"""
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    bullet_layout = prs.slide_layouts[1]

    # 封面
    s = prs.slides.add_slide(title_slide_layout)
    s.shapes.title.text = REPORT_TITLE
    s.placeholders[1].text = "灵犀演示 · Phase 1 · 季度汇报场景"

    # 业绩摘要
    s = prs.slides.add_slide(bullet_layout)
    s.shapes.title.text = "Q1 业绩摘要"
    body = s.shapes.placeholders[1].text_frame
    body.text = REPORT_PARAGRAPHS[0]
    for p in REPORT_PARAGRAPHS[1:4]:
        para = body.add_paragraph()
        para.text = p

    # AI 产品线
    s = prs.slides.add_slide(bullet_layout)
    s.shapes.title.text = "AI 产品线进展"
    body = s.shapes.placeholders[1].text_frame
    body.text = REPORT_PARAGRAPHS[2]
    for p in [
        "5 模块独立 demo 全跑通（文件 / 顾问 / 模板 / 预览 / 输出）",
        "AI 顾问交互准确率 ≥ 90%",
        "100M 文件导入成功率 ≥ 99%",
    ]:
        para = body.add_paragraph()
        para.text = p

    # 海外市场
    s = prs.slides.add_slide(bullet_layout)
    s.shapes.title.text = "海外市场"
    body = s.shapes.placeholders[1].text_frame
    body.text = REPORT_PARAGRAPHS[3]
    for p in [
        "新加坡 POC 1（航空维修 SaaS）",
        "印尼 POC 1（航运代理）",
        "马来西亚 POC 1（制造业 ERP）",
    ]:
        para = body.add_paragraph()
        para.text = p

    # Q2 计划
    s = prs.slides.add_slide(bullet_layout)
    s.shapes.title.text = "Q2 计划"
    body = s.shapes.placeholders[1].text_frame
    body.text = REPORT_PARAGRAPHS[6]
    for p in [
        "付费版上线（ARR 200 万）",
        "行业模板 5 套（航材 / 航运 / 制造）",
        "海外 POC 转化 2 单",
    ]:
        para = body.add_paragraph()
        para.text = p

    out = OUT / "sample.pptx"
    prs.save(out)
    print(f"  [OK] {out.name} ({out.stat().st_size} bytes)")


def make_md():
    """Markdown 样本"""
    out = OUT / "sample.md"
    lines = [
        f"# {REPORT_TITLE}",
        "",
        "## 业绩摘要",
        "",
    ]
    for p in REPORT_PARAGRAPHS:
        lines.append(f"- {p}")
    lines += [
        "",
        "## 关键数字",
        "",
        "| 指标 | Q1 | Q4 2025 | 同比 |",
        "| --- | --- | --- | --- |",
        "| 营收(万) | 1280 | 1040 | +23% |",
        "| 新增客户 | 14 | 9 | +56% |",
        "| AI 产品线占比 | 41% | 28% | +13pp |",
        "",
        "## 关键词",
        "",
        " ".join(f"`{t}`" for t in TAGS),
        "",
        "---",
        "",
        "_本文件由 T-1.1 testdata generator 生成。_",
        "",
    ]
    out.write_text("\n".join(lines), encoding="utf-8")
    print(f"  [OK] {out.name} ({out.stat().st_size} bytes)")


def make_image(fmt: str, name: str):
    """图片样本（800x600 + CJK 文字水印）"""
    img = Image.new("RGB" if fmt == "jpg" else "RGBA", (800, 600), color=(245, 245, 250) if fmt == "jpg" else (245, 245, 250, 255))
    draw = ImageDraw.Draw(img)
    # 找一个能写中文的字体（macOS PingFang / Heiti；找不到就用 default）
    font = None
    for path in [
        "/System/Library/Fonts/PingFang.ttc",
        "/System/Library/Fonts/STHeiti Medium.ttc",
        "/System/Library/Fonts/Helvetica.ttc",
        "/Library/Fonts/Arial Unicode.ttf",
    ]:
        if os.path.exists(path):
            try:
                font = ImageFont.truetype(path, 36)
                break
            except Exception:
                continue
    if font is None:
        font = ImageFont.load_default()
    # 标题
    draw.text((40, 40), REPORT_TITLE, fill=(20, 20, 80), font=font)
    # 摘要 4 行
    for i, p in enumerate(REPORT_PARAGRAPHS[:4]):
        draw.text((60, 130 + i * 50), f"• {p[:30]}…", fill=(60, 60, 60), font=font)
    # 关键词
    draw.text((40, 460), "标签: " + " ".join(f"#{t}" for t in TAGS[:5]), fill=(120, 60, 60), font=font)
    # 标识
    draw.text((40, 540), f"T-1.1 testdata · {fmt.upper()} sample · 800x600", fill=(160, 160, 160), font=font)

    out = OUT / name
    if fmt == "jpg":
        img.save(out, "JPEG", quality=85)
    else:
        img.save(out, "PNG")
    print(f"  [OK] {out.name} ({out.stat().st_size} bytes)")


def write_readme():
    """testdata 索引"""
    out = OUT / "README.md"
    out.write_text(
        """# T-1.1 测试数据

7 格式样本 + 1 个大 PDF（>50MB，用于 100M 内成功率压测）。

| 文件 | 大小 | 用途 |
| --- | --- | --- |
| sample.docx | ~38KB | Word 文档样本 |
| sample.pdf | ~8KB | PDF 文档样本（3 页） |
| sample.xlsx | ~6KB | Excel 表格样本 |
| sample.pptx | ~50KB | PowerPoint 样本（5 页） |
| sample.md | ~2KB | Markdown 笔记样本 |
| sample.jpg | ~80KB | JPG 图片样本（800x600） |
| sample.png | ~50KB | PNG 图片样本（800x600） |
| large_50mb.pdf | ~55MB | 100M 内压测样本（PRD ≥99% 成功率验证） |

生成：`python3 scripts/generate_testdata.py`

灵犀演示 · Phase 1 · T-1.1
""",
        encoding="utf-8",
    )
    print(f"  [OK] {out.name}")


def main():
    print(f"Generating testdata in {OUT} ...")
    make_docx()
    make_pdf(pages=3)
    make_xlsx()
    make_pptx()
    make_md()
    make_image("jpg", "sample.jpg")
    make_image("png", "sample.png")
    # 60MB 大 PDF（PRD 上限 100M，留余量）
    make_pdf(pages=60, target_size_bytes=60 * 1024 * 1024, name="large_50mb.pdf")
    write_readme()
    print("\nDone. 8 files + README.md.")


if __name__ == "__main__":
    main()
