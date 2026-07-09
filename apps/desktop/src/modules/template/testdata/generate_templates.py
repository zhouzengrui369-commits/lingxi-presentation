"""生成 3 套测试 PPTX 模板。

每套 ≥10 页，3 种截然不同的风格（商务深色 / 学术浅色 / 创意渐变）。

用法:
    /usr/bin/python3 apps/desktop/src/modules/template/testdata/generate_templates.py

输出:
    apps/desktop/testdata/templates/business-dark.pptx
    apps/desktop/testdata/templates/academic-light.pptx
    apps/desktop/testdata/templates/creative-gradient.pptx
"""
from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

OUT_DIR = Path(__file__).resolve().parent / "templates"
OUT_DIR.mkdir(parents=True, exist_ok=True)

# 同时写到 spec 规定的路径（apps/desktop/testdata/templates/）
SPEC_OUT_DIR = Path(__file__).resolve().parents[4] / "testdata" / "templates"
SPEC_OUT_DIR.mkdir(parents=True, exist_ok=True)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _add_title(slide, text, *, color="FFFFFF", size=44, bold=True, x=Inches(0.6), y=Inches(0.6),
               w=Inches(12), h=Inches(1.2), font="Microsoft YaHei"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor.from_string(color)
    run.font.name = font
    return tb


def _add_body(slide, text, *, color="FFFFFF", size=18, x=Inches(0.6), y=Inches(2.0),
              w=Inches(12), h=Inches(4.5), bullet=False, font="Microsoft YaHei"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = (f"• {line}" if bullet else line)
        run.font.size = Pt(size)
        run.font.color.rgb = RGBColor.from_string(color)
        run.font.name = font
    return tb


# Per-template font shortcuts (avoid 60+ kwargs passes)
def _add_title_bd(slide, text, **kw):  # business-dark
    kw.setdefault("font", "Microsoft YaHei")
    return _add_title(slide, text, **kw)


def _add_body_bd(slide, text, **kw):
    kw.setdefault("font", "Microsoft YaHei")
    return _add_body(slide, text, **kw)


def _add_title_ac(slide, text, **kw):  # academic-light
    kw.setdefault("font", "Source Han Sans CN")
    return _add_title(slide, text, **kw)


def _add_body_ac(slide, text, **kw):
    kw.setdefault("font", "Source Han Sans CN")
    return _add_body(slide, text, **kw)


def _add_title_cg(slide, text, **kw):  # creative-gradient
    kw.setdefault("font", "PingFang SC")
    return _add_title(slide, text, **kw)


def _add_body_cg(slide, text, **kw):
    kw.setdefault("font", "PingFang SC")
    return _add_body(slide, text, **kw)


def _set_bg(slide, hex_color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(hex_color)


def _add_accent_bar(slide, *, color, top=Inches(5.6), height=Inches(0.08), x=Inches(0.6), w=Inches(12.1)):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, top, w, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor.from_string(color)
    bar.line.fill.background()


def _add_block(slide, *, color, x, y, w, h):
    """装饰色块（创意渐变模板用）"""
    block = slide.shapes.add_shape(MSO_SHAPE.OVAL if "oval" in str(color) else MSO_SHAPE.RECTANGLE, x, y, w, h)
    block.fill.solid()
    block.fill.fore_color.rgb = RGBColor.from_string(color)
    block.line.fill.background()


# ---- 1. business-dark ----

def build_business_dark() -> Path:
    """商务深色模板: 深蓝底 + 金色重点 + 衬线标题"""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    # 1. title slide
    s = prs.slides.add_slide(blank)
    _set_bg(s, "0E1421")
    _add_title_bd(s, "2026 Q1 业务回顾", color="C9A227", size=60, y=Inches(2.5))
    _add_body_bd(s, "董事会汇报 · 仅限内部", color="E8E8E8", size=24, y=Inches(4.0), h=Inches(1.2))
    _add_accent_bar(s, color="C9A227")

    # 2. section: 业绩概览
    s = prs.slides.add_slide(blank)
    _set_bg(s, "0E1421")
    _add_title_bd(s, "01 业绩概览", color="C9A227", size=36, y=Inches(0.6))
    _add_body_bd(s, "营收达成率 108% · 同比增长 23%", color="E8E8E8", size=22, y=Inches(2.5))

    # 3. content: 关键数字
    s = prs.slides.add_slide(blank)
    _set_bg(s, "1A1A2E")
    _add_title_bd(s, "关键数字", color="C9A227", size=32, y=Inches(0.5))
    bullets = "营收 1.42 亿\n毛利 5,800 万\n净利 2,310 万\nARPU 同比 +18%"
    _add_body_bd(s, bullets, color="FFFFFF", size=22, y=Inches(2.0), bullet=True)

    # 4. content: 区域表现
    s = prs.slides.add_slide(blank)
    _set_bg(s, "1A1A2E")
    _add_title_bd(s, "区域表现", color="C9A227", size=32, y=Inches(0.5))
    bullets = "华东: 5,400 万 (+25%)\n华南: 3,800 万 (+18%)\n华北: 3,200 万 (+30%)\n西部: 1,800 万 (+15%)"
    _add_body_bd(s, bullets, color="FFFFFF", size=20, y=Inches(2.0), bullet=True)

    # 5. two-column: 客户增长
    s = prs.slides.add_slide(blank)
    _set_bg(s, "0E1421")
    _add_title_bd(s, "客户增长", color="C9A227", size=32, y=Inches(0.5))
    _add_body_bd(s, "新增企业客户\n42 家", color="FFFFFF", size=28, x=Inches(0.6), y=Inches(2.5), w=Inches(6))
    _add_body_bd(s, "续约率\n91%", color="C9A227", size=28, x=Inches(6.8), y=Inches(2.5), w=Inches(6))

    # 6. section: 产品线
    s = prs.slides.add_slide(blank)
    _set_bg(s, "0E1421")
    _add_title_bd(s, "02 产品线", color="C9A227", size=36, y=Inches(0.6))

    # 7. content: 产品矩阵
    s = prs.slides.add_slide(blank)
    _set_bg(s, "1A1A2E")
    _add_title_bd(s, "产品矩阵", color="C9A227", size=32, y=Inches(0.5))
    _add_body_bd(s, "旗舰版 · 增长版 · 入门版 · 定制版", color="FFFFFF", size=24, y=Inches(2.5))

    # 8. quote: 客户原声
    s = prs.slides.add_slide(blank)
    _set_bg(s, "0E1421")
    _add_body_bd(s, '"灵犀演示让我们的季度汇报从 3 天压缩到 4 小时。"', color="C9A227", size=28,
              x=Inches(1.5), y=Inches(2.8), w=Inches(10.3), h=Inches(2))
    _add_body_bd(s, "— 张总监 · 某制造业上市公司", color="FFFFFF", size=18, y=Inches(5.0))

    # 9. content: 风险与挑战
    s = prs.slides.add_slide(blank)
    _set_bg(s, "1A1A2E")
    _add_title_bd(s, "风险与挑战", color="C9A227", size=32, y=Inches(0.5))
    bullets = "原材料成本上涨 12%\n汇率波动影响海外营收\n新业务线尚在投入期\n监管政策变化"
    _add_body_bd(s, bullets, color="FFFFFF", size=22, y=Inches(2.0), bullet=True)

    # 10. summary
    s = prs.slides.add_slide(blank)
    _set_bg(s, "0E1421")
    _add_title_bd(s, "Q2 展望", color="C9A227", size=36, y=Inches(0.6))
    bullets = "聚焦核心产品迭代\n扩大华东华南覆盖\n启动海外试点\n目标营收 1.6 亿"
    _add_body_bd(s, bullets, color="FFFFFF", size=24, y=Inches(2.5), bullet=True)
    _add_accent_bar(s, color="C9A227")

    out = OUT_DIR / "business-dark.pptx"
    prs.save(str(out))
    return out


# ---- 2. academic-light ----

def build_academic_light() -> Path:
    """学术浅色模板: 白底 + 蓝灰重点 + 无衬线字体"""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    # 1. title
    s = prs.slides.add_slide(blank)
    _set_bg(s, "FFFFFF")
    _add_title_ac(s, "基于大语言模型的\n代码生成综述", color="1E40AF", size=48, y=Inches(2.0), h=Inches(2.5))
    _add_body_ac(s, "论文答辩 · 2026 春季", color="475569", size=20, y=Inches(5.0))
    _add_accent_bar(s, color="3B82F6", top=Inches(6.5), height=Inches(0.05))

    # 2. section: 研究背景
    s = prs.slides.add_slide(blank)
    _set_bg(s, "FFFFFF")
    _add_title_ac(s, "01 研究背景", color="1E40AF", size=32, y=Inches(0.6))
    _add_body_ac(s, "LLM 在代码生成领域取得显著进展", color="334155", size=22, y=Inches(2.5))

    # 3. content: 研究问题
    s = prs.slides.add_slide(blank)
    _set_bg(s, "F8FAFC")
    _add_title_ac(s, "研究问题", color="1E40AF", size=28, y=Inches(0.5))
    bullets = "现有方法在长上下文生成上的局限？\n不同模型在多语言场景的差异？\n代码质量评估指标的局限性？"
    _add_body_ac(s, bullets, color="334155", size=20, y=Inches(2.0), bullet=True)

    # 4. content: 文献综述
    s = prs.slides.add_slide(blank)
    _set_bg(s, "F8FAFC")
    _add_title_ac(s, "文献综述", color="1E40AF", size=28, y=Inches(0.5))
    _add_body_ac(s, "我们综述了 2018-2026 年间 87 篇相关论文，覆盖 5 大研究主题。",
              color="334155", size=20, y=Inches(2.0))

    # 5. two-column: 方法对比
    s = prs.slides.add_slide(blank)
    _set_bg(s, "FFFFFF")
    _add_title_ac(s, "方法对比", color="1E40AF", size=28, y=Inches(0.5))
    _add_body_ac(s, "基于 Prompt 的方法\n直观但难优化", color="334155", size=22,
              x=Inches(0.6), y=Inches(2.5), w=Inches(6))
    _add_body_ac(s, "基于微调的方法\n效果好但成本高", color="334155", size=22,
              x=Inches(6.8), y=Inches(2.5), w=Inches(6))

    # 6. section: 方法
    s = prs.slides.add_slide(blank)
    _set_bg(s, "FFFFFF")
    _add_title_ac(s, "02 我们的方法", color="1E40AF", size=32, y=Inches(0.6))

    # 7. content: 实验设置
    s = prs.slides.add_slide(blank)
    _set_bg(s, "F8FAFC")
    _add_title_ac(s, "实验设置", color="1E40AF", size=28, y=Inches(0.5))
    bullets = "数据集: HumanEval + MBPP + 自建\n基线: Codex / CodeLlama / GPT-4\n评估: pass@k / BLEU / CodeBLEU"
    _add_body_ac(s, bullets, color="334155", size=20, y=Inches(2.0), bullet=True)

    # 8. content: 实验结果
    s = prs.slides.add_slide(blank)
    _set_bg(s, "F8FAFC")
    _add_title_ac(s, "实验结果", color="1E40AF", size=28, y=Inches(0.5))
    _add_body_ac(s, "我们的方法在 HumanEval 上 pass@1 达到 78.5%，比基线最优高 4.2%。",
              color="334155", size=20, y=Inches(2.0))

    # 9. quote: 关键发现
    s = prs.slides.add_slide(blank)
    _set_bg(s, "FFFFFF")
    _add_body_ac(s, '"在长上下文场景下，检索增强能将准确率提升 12%。"', color="1E40AF",
              size=28, x=Inches(1.5), y=Inches(2.8), w=Inches(10.3), h=Inches(2))

    # 10. summary
    s = prs.slides.add_slide(blank)
    _set_bg(s, "FFFFFF")
    _add_title_ac(s, "结论与展望", color="1E40AF", size=32, y=Inches(0.6))
    bullets = "提出检索增强生成框架\n在 5 个数据集上验证\n未来: 多模态代码理解"
    _add_body_ac(s, bullets, color="334155", size=24, y=Inches(2.5), bullet=True)
    _add_accent_bar(s, color="3B82F6", top=Inches(6.5), height=Inches(0.05))

    out = OUT_DIR / "academic-light.pptx"
    prs.save(str(out))
    return out


# ---- 3. creative-gradient ----

def build_creative_gradient() -> Path:
    """创意渐变模板: 多彩底 + 圆角色块 + 现代装饰"""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    # 1. title: 多色块叠加
    s = prs.slides.add_slide(blank)
    _set_bg(s, "FAF5FF")
    _add_block(s, color="A855F7", x=Inches(-1), y=Inches(-1), w=Inches(8), h=Inches(8))
    _add_block(s, color="EC4899", x=Inches(5), y=Inches(2), w=Inches(7), h=Inches(6))
    _add_block(s, color="F59E0B", x=Inches(10), y=Inches(5), w=Inches(5), h=Inches(4))
    _add_title_cg(s, "品牌升级\n从 0 到 1", color="FFFFFF", size=56, y=Inches(2.5), h=Inches(2.5))

    # 2. section
    s = prs.slides.add_slide(blank)
    _set_bg(s, "FCE7F3")
    _add_title_cg(s, "01 品牌故事", color="831843", size=40, y=Inches(0.6))
    _add_body_cg(s, "从工坊到全球", color="9D174D", size=24, y=Inches(2.5))

    # 3. content: 三个色块表示三个支柱
    s = prs.slides.add_slide(blank)
    _set_bg(s, "FFFFFF")
    _add_title_cg(s, "三大支柱", color="581C87", size=32, y=Inches(0.5))
    _add_block(s, color="A855F7", x=Inches(0.8), y=Inches(2.5), w=Inches(3.5), h=Inches(3.5))
    _add_block(s, color="EC4899", x=Inches(4.9), y=Inches(2.5), w=Inches(3.5), h=Inches(3.5))
    _add_block(s, color="F59E0B", x=Inches(9), y=Inches(2.5), w=Inches(3.5), h=Inches(3.5))
    _add_body_cg(s, "创新", color="FFFFFF", size=28, x=Inches(0.8), y=Inches(3.7), w=Inches(3.5))
    _add_body_cg(s, "品质", color="FFFFFF", size=28, x=Inches(4.9), y=Inches(3.7), w=Inches(3.5))
    _add_body_cg(s, "社群", color="FFFFFF", size=28, x=Inches(9), y=Inches(3.7), w=Inches(3.5))

    # 4. content: 数据
    s = prs.slides.add_slide(blank)
    _set_bg(s, "FDF4FF")
    _add_title_cg(s, "增长数据", color="581C87", size=32, y=Inches(0.5))
    bullets = "月活: 50 万 → 230 万\n转化率: 2.1% → 5.8%\nNPS: 42 → 67\n客单价: ¥299 → ¥459"
    _add_body_cg(s, bullets, color="581C87", size=22, y=Inches(2.0), bullet=True)

    # 5. two-column: 用户故事
    s = prs.slides.add_slide(blank)
    _set_bg(s, "FFFFFF")
    _add_title_cg(s, "用户故事", color="581C87", size=32, y=Inches(0.5))
    _add_block(s, color="A855F7", x=Inches(0.6), y=Inches(2.2), w=Inches(5.8), h=Inches(4))
    _add_block(s, color="EC4899", x=Inches(6.9), y=Inches(2.2), w=Inches(5.8), h=Inches(4))
    _add_body_cg(s, '"我每天早上都用它来开始创作。"', color="FFFFFF", size=18,
              x=Inches(0.8), y=Inches(3.0), w=Inches(5.4), h=Inches(2.5))
    _add_body_cg(s, '"灵犀让品牌升级像游戏一样有趣。"', color="FFFFFF", size=18,
              x=Inches(7.1), y=Inches(3.0), w=Inches(5.4), h=Inches(2.5))

    # 6. section
    s = prs.slides.add_slide(blank)
    _set_bg(s, "FAE8FF")
    _add_title_cg(s, "02 未来路线", color="831843", size=40, y=Inches(0.6))

    # 7. content: 路线图
    s = prs.slides.add_slide(blank)
    _set_bg(s, "FFFFFF")
    _add_title_cg(s, "2026 路线图", color="581C87", size=32, y=Inches(0.5))
    bullets = "Q1: AI 助手上线\nQ2: 移动端 App\nQ3: 全球扩张\nQ4: IPO 准备"
    _add_body_cg(s, bullets, color="581C87", size=24, y=Inches(2.0), bullet=True)

    # 8. quote: 愿景
    s = prs.slides.add_slide(blank)
    _set_bg(s, "FDF4FF")
    _add_block(s, color="F59E0B", x=Inches(5), y=Inches(2), w=Inches(3.3), h=Inches(3.3))
    _add_body_cg(s, '"让每个人都能成为创作者。"', color="581C87", size=32,
              x=Inches(1.5), y=Inches(3.0), w=Inches(10.3), h=Inches(1.5))

    # 9. content: 团队
    s = prs.slides.add_slide(blank)
    _set_bg(s, "FFFFFF")
    _add_title_cg(s, "团队介绍", color="581C87", size=32, y=Inches(0.5))
    _add_body_cg(s, "60+ 同事 · 来自 12 个国家 · 平均年龄 28", color="831843", size=22, y=Inches(2.5))

    # 10. summary
    s = prs.slides.add_slide(blank)
    _set_bg(s, "FAF5FF")
    _add_block(s, color="A855F7", x=Inches(-1), y=Inches(5), w=Inches(14), h=Inches(3))
    _add_block(s, color="EC4899", x=Inches(-1), y=Inches(6.5), w=Inches(14), h=Inches(2))
    _add_title_cg(s, "谢谢！", color="FFFFFF", size=72, y=Inches(2.5))

    out = OUT_DIR / "creative-gradient.pptx"
    prs.save(str(out))
    return out


def main() -> int:
    out1 = build_business_dark()
    print(f"[ok] {out1} ({out1.stat().st_size // 1024} KB)")
    out2 = build_academic_light()
    print(f"[ok] {out2} ({out2.stat().st_size // 1024} KB)")
    out3 = build_creative_gradient()
    print(f"[ok] {out3} ({out3.stat().st_size // 1024} KB)")

    # 复制到 spec 规定路径 apps/desktop/testdata/templates/，方便 PM 验收 + 后续 Phase 2/3 集成
    import shutil
    for src in [out1, out2, out3]:
        dst = SPEC_OUT_DIR / src.name
        shutil.copy2(src, dst)
        print(f"[ok] copied → {dst}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
